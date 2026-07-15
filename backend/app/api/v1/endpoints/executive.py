from datetime import datetime
from uuid import UUID

from fastapi import APIRouter, Depends, HTTPException, Query, status
from sqlalchemy.orm import Session

from app.api.v1.deps import get_current_active_user, get_db
from app.models.membership import Membership
from app.models.user import User
from app.schemas.executive import (
    CEOGenerationRequest,
    CEOReportResponse,
    COOGenerationRequest,
    COOReportResponse,
    CTOGenerationRequest,
    CTOReportResponse,
)
from app.services.executive.ceo import ceo_service
from app.services.executive.coo import coo_service
from app.services.executive.cto import cto_service

from app.repositories.executive import ceo_report_repo, cto_report_repo, coo_report_repo
from app.models.executive import CEOReport, CTOReport, COOReport

router = APIRouter()


def check_workspace_access(db: Session, user_id: UUID, workspace_id: UUID):
    """
    Checks workspace membership for the user.
    """
    membership = (
        db.query(Membership)
        .filter(Membership.user_id == user_id, Membership.workspace_id == workspace_id)
        .first()
    )
    if not membership:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="User does not have access permissions for this workspace.",
        )


from app.models.workspace import Workspace

def get_report_securely(
    db: Session,
    report_id: UUID,
    user_id: UUID,
):
    """
    Retrieve report strictly verifying active workspace membership.
    Derives allowed workspaces directly from the user rather than trusting client parameters.
    """
    # Join Workspace to ensure workspace soft-deletion and archived state are enforced
    query = db.query(Membership.workspace_id).join(
        Workspace, Workspace.id == Membership.workspace_id
    ).filter(
        Membership.user_id == user_id,
        Membership.deleted_at.is_(None),
        Workspace.deleted_at.is_(None),
        Workspace.archived == False,
    )

    workspaces = list({w[0] for w in query.all()})


    if not workspaces:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Report not found",
        )

    # Use repository classes consistently to query reports securely (avoiding direct db.query)
    ceo = ceo_report_repo.get_by_id_and_workspaces(db, id=report_id, workspace_ids=workspaces)
    if ceo:
        return ceo, "ceo"

    cto = cto_report_repo.get_by_id_and_workspaces(db, id=report_id, workspace_ids=workspaces)
    if cto:
        return cto, "cto"

    coo = coo_report_repo.get_by_id_and_workspaces(db, id=report_id, workspace_ids=workspaces)
    if coo:
        return coo, "coo"

    raise HTTPException(
        status_code=status.HTTP_404_NOT_FOUND,
        detail="Report not found",
    )




@router.post("/ceo/report", response_model=CEOReportResponse)
def generate_ceo_report(
    req: CEOGenerationRequest,
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    try:
        return ceo_service.generate_ceo_report(
            db=db,
            workspace_id=workspace_id,
            product_idea=req.product_idea,
            target_industry=req.target_industry or "Tech",
            competitors=req.competitors,
            budget=req.budget or 100000.0,
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate CEO report: {str(e)}",
        )


@router.get("/ceo/reports", response_model=list[CEOReportResponse])
def get_ceo_reports(
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    return ceo_service.get_ceo_reports(db, workspace_id)


@router.post("/cto/report", response_model=CTOReportResponse)
def generate_cto_report(
    req: CTOGenerationRequest,
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    try:
        return cto_service.generate_cto_report(
            db=db,
            workspace_id=workspace_id,
            product_spec=req.product_spec,
            preferred_cloud=req.preferred_cloud or "AWS",
            compliance_needs=req.compliance_needs,
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate CTO report: {str(e)}",
        )


@router.get("/cto/reports", response_model=list[CTOReportResponse])
def get_cto_reports(
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    return cto_service.get_cto_reports(db, workspace_id)


@router.post("/coo/report", response_model=COOReportResponse)
def generate_coo_report(
    req: COOGenerationRequest,
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    try:
        return coo_service.generate_coo_report(
            db=db,
            workspace_id=workspace_id,
            sprint_name=req.sprint_name,
            team_members=req.team_members,
            total_budget=req.total_budget or 50000.0,
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate COO report: {str(e)}",
        )


@router.get("/coo/reports", response_model=list[COOReportResponse])
def get_coo_reports(
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    check_workspace_access(db, current_user.id, workspace_id)
    return coo_service.get_coo_reports(db, workspace_id)


@router.get("/export/pdf")
def export_pdf_report(
    report_id: UUID,
    report_type: str = Query(..., pattern="^(ceo|cto|coo)$"),
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    """
    Simulated executive business report PDF exporter.
    """
    check_workspace_access(db, current_user.id, workspace_id)
    return {
        "success": True,
        "message": f"Successfully compiled and generated PDF export for {report_type} report: {report_id}.",
        "download_url": f"/api/v1/executive/download/pdf/{report_id}",
    }


@router.get("/export/ppt")
def export_ppt_report(
    report_id: UUID,
    report_type: str = Query(..., pattern="^(ceo|cto|coo)$"),
    workspace_id: UUID = Query(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    """
    Simulated executive presentation PPT exporter.
    """
    check_workspace_access(db, current_user.id, workspace_id)
    return {
        "success": True,
        "message": f"Successfully compiled and generated PPT slides deck export for {report_type} report: {report_id}.",
        "download_url": f"/api/v1/executive/download/ppt/{report_id}",
    }


@router.get("/download/pdf/{report_id}")
def download_pdf_report(
    report_id: UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    """
    Serve the generated PDF report content as a real binary PDF document.
    """
    import io
    from fastapi.responses import StreamingResponse
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor

    from app.models.executive import CEOReport, COOReport, CTOReport

    report, report_type = get_report_securely(
        db, report_id=report_id, user_id=current_user.id
    )


    title = "Executive Advisor Report"
    content = "Executive Summary:\n\n"
    if report_type == "ceo":
        ceo = report
        title = ceo.title
        content += f"Active Idea: {ceo.portfolio_data.get('active_idea', '')}\n"
        content += f"Target Vertical: {ceo.portfolio_data.get('target_vertical', '')}\n\n"
        content += "SWOT Analysis:\n"
        content += f"{ceo.strategy_data.get('swot', {}).get('swot_markdown', '')}\n\n"
        content += "PESTLE Analysis:\n"
        content += f"{ceo.strategy_data.get('pestle', {}).get('pestle_markdown', '')}\n\n"
        content += f"Financial Forecast:\n"
        content += f"{ceo.financials.get('forecast', {}).get('forecast_markdown', '')}\n\n"
    elif report_type == "cto":
        cto = report
        title = cto.title
        content += f"Tech Stack Spec: {cto.spec_data.get('product_spec', '')}\n"
        content += f"Preferred Cloud: {cto.architecture_data.get('preferred_cloud', '')}\n\n"
        content += "System Design:\n"
        content += f"{cto.architecture_data.get('design_markdown', '')}\n\n"
        content += "API Specifications:\n"
        content += f"{cto.spec_data.get('api_markdown', '')}\n\n"
    elif report_type == "coo":
        coo = report
        title = coo.title
        content += f"Sprint Name: {coo.operations_data.get('sprint_name', '')}\n"
        content += "Sprint Backlog & Tasks:\n"
        content += f"{coo.operations_data.get('tasks_markdown', '')}\n\n"
        content += "Team Resources:\n"
        content += f"{coo.resource_data.get('team_members', '')}\n\n"


    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=20,
        leading=24,
        textColor=HexColor('#1E293B'),
        spaceAfter=15
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontSize=10,
        leading=14,
        textColor=HexColor('#334155'),
        spaceAfter=6
    )

    story = []
    story.append(Paragraph(title.upper(), title_style))
    story.append(Spacer(1, 10))
    
    for line in content.split("\n"):
        line_clean = line.strip().replace("<", "&lt;").replace(">", "&gt;")
        if line_clean:
            story.append(Paragraph(line_clean, body_style))
        else:
            story.append(Spacer(1, 5))
            
    doc.build(story)
    buf.seek(0)
    
    filename = f"{title.replace(' ', '_').lower()}.pdf"
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@router.get("/download/ppt/{report_id}")
def download_ppt_report(
    report_id: UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user),
):
    """
    Serve the generated PPT slides report content as a real binary PPTX file.
    """
    import io
    from fastapi.responses import StreamingResponse
    from pptx import Presentation

    from app.models.executive import CEOReport, COOReport, CTOReport

    report, report_type = get_report_securely(
        db, report_id=report_id, user_id=current_user.id
    )


    title = "Executive Presentation"
    slides = []
    if report_type == "ceo":
        ceo = report
        title = ceo.title
        slides = [
            (
                "Executive Presentation: " + ceo.title,
                f"Active Idea: {ceo.portfolio_data.get('active_idea', '')}\nTarget Vertical: {ceo.portfolio_data.get('target_vertical', '')}"
            ),
            (
                "SWOT Analysis",
                f"{ceo.strategy_data.get('swot', {}).get('swot_markdown', '')}"
            ),
            (
                "PESTLE Analysis",
                f"{ceo.strategy_data.get('pestle', {}).get('pestle_markdown', '')}"
            ),
            (
                "Financial Forecast & Budget",
                f"Initial Budget: ${ceo.financials.get('initial_budget', 0)}\nForecast Details: {ceo.financials.get('forecast', {}).get('forecast_markdown', '')}"
            ),
        ]
    elif report_type == "cto":
        cto = report
        title = cto.title
        slides = [
            (
                "Architecture Overview",
                f"Technical Specifications for: {cto.title}"
            ),
            (
                "System Design & Strategy",
                f"{cto.architecture_data.get('design_markdown', '')}"
            ),
            (
                "API & Schema Specs",
                f"{cto.spec_data.get('api_markdown', '')}"
            ),
        ]
    elif report_type == "coo":
        coo = report
        title = coo.title
        slides = [
            (
                "Operations Backlog",
                f"Project Backlog and Sprint for: {coo.title}"
            ),
            (
                "Sprint Board Status",
                f"{coo.operations_data.get('tasks_markdown', '')}"
            ),
            (
                "Team Resource Allocation",
                f"{coo.resource_data.get('team_members', '')}"
            ),
        ]


    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = "Executive Advisory Presentation\nAI ProductOS Boardroom"
    
    # Content Slides
    for slide_title, slide_content in slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_title
        body_shape.text = slide_content

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    filename = f"{title.replace(' ', '_').lower()}_presentation.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )
